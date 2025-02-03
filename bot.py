import os
from flask import Flask, request, render_template, redirect, url_for
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

app = Flask(__name__)

# Global variables to hold PPT text segments and the index
ppt_file = "rr_01_artificial_intelligence.pptx"  # Update with your PPT file path
slides_text = []
vectorizer = None
tfidf_matrix = None

def extract_text_from_ppt(file_path):
    """Extracts text from each slide in the PowerPoint file."""
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text = "\n".join(slide_text)
        slides_text.append(text)
    return slides_text

def build_index(text_segments):
    """Builds a TF-IDF vectorizer on the text segments."""
    vectorizer = TfidfVectorizer(stop_words='english')
    tfidf_matrix = vectorizer.fit_transform(text_segments)
    return vectorizer, tfidf_matrix

def find_relevant_segment(question, vectorizer, tfidf_matrix, text_segments, top_n=1):
    """Given a question, returns the top_n most similar text segments."""
    question_vec = vectorizer.transform([question])
    similarities = cosine_similarity(question_vec, tfidf_matrix).flatten()
    top_indices = similarities.argsort()[-top_n:][::-1]
    return [(idx, text_segments[idx], similarities[idx]) for idx in top_indices]

@app.before_request
def initialize():
    """Initialize by extracting text and building the index."""
    global slides_text, vectorizer, tfidf_matrix
    if not os.path.exists(ppt_file):
        raise FileNotFoundError(f"File not found: {ppt_file}")
    slides_text = extract_text_from_ppt(ppt_file)
    vectorizer, tfidf_matrix = build_index(slides_text)
    print("Initialization complete: PPT text extracted and index built.")

@app.route("/", methods=["GET", "POST"])
def index():
    result = None
    if request.method == "POST":
        question = request.form.get("question", "")
        if question:
            results = find_relevant_segment(question, vectorizer, tfidf_matrix, slides_text)
            # Taking the top result
            if results:
                idx, text, score = results[0]
                result = {
                    "slide_number": idx + 1,
                    "content": text,
                    "score": f"{score:.2f}"
                }
    return render_template("index.html", result=result)

if __name__ == "__main__":
    app.run(debug=True)
